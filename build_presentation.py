"""
Multi-Agent AI Academic Assistant — Final Presentation
Dark navy + teal + green, visually rich, 13 slides (16:9)
Run: python build_presentation.py
Output: presentation_final.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0B, 0x14, 0x2A)
CARD    = RGBColor(0x14, 0x22, 0x3E)
TEAL    = RGBColor(0x00, 0xC9, 0xB1)
GREEN   = RGBColor(0x22, 0xD3, 0x6A)
AMBER   = RGBColor(0xFB, 0xBF, 0x24)
PURPLE  = RGBColor(0xA7, 0x8B, 0xFA)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x8A, 0x9B, 0xB4)
LIGHT   = RGBColor(0xE2, 0xE8, 0xF0)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────
def new_slide():
    s = prs.slides.add_slide(blank)
    f = s.background.fill
    f.solid(); f.fore_color.rgb = NAVY
    return s

def rect(slide, l, t, w, h, fill, border_color=None, border_pt=0):
    sp = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if border_color:
        sp.line.color.rgb = border_color
        sp.line.width = Pt(border_pt)
    else:
        sp.line.fill.background()
    return sp

def label(slide, text, l, t, w, h, size=16, bold=False,
          color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def pill(slide, text, l, t, color, size=11):
    w = len(text) * 0.088 + 0.28
    rect(slide, l, t, w, 0.3, color)
    label(slide, text, l+0.1, t+0.04, w-0.1, 0.24,
          size=size, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

def slide_header(slide, title, tag="", tag_color=TEAL):
    rect(slide, 0, 0, 13.33, 0.06, TEAL)
    if tag:
        pill(slide, tag, 0.5, 0.18, tag_color)
        label(slide, title, 0.5, 0.52, 12.3, 0.65, size=28, bold=True, color=WHITE)
    else:
        label(slide, title, 0.5, 0.18, 12.3, 0.72, size=28, bold=True, color=WHITE)
    rect(slide, 0.5, 1.1, 12.3, 0.025, CARD)

def arrow_h(slide, l, t, w, color=MUTED):
    """Horizontal arrow: line + triangle head."""
    rect(slide, l, t + 0.045, w - 0.12, 0.025, color)
    # arrowhead triangle
    tri = slide.shapes.add_shape(5, Inches(l + w - 0.14), Inches(t),
                                  Inches(0.14), Inches(0.115))
    tri.fill.solid(); tri.fill.fore_color.rgb = color
    tri.line.fill.background()

def arrow_v(slide, l, t, h, color=MUTED):
    """Vertical arrow: line + triangle head."""
    rect(slide, l + 0.045, t, 0.025, h - 0.12, color)
    tri = slide.shapes.add_shape(6, Inches(l), Inches(t + h - 0.14),
                                  Inches(0.115), Inches(0.14))
    tri.fill.solid(); tri.fill.fore_color.rgb = color
    tri.line.fill.background()

def node(slide, l, t, w, h, title, subtitle, color, icon=""):
    """A labelled diagram node (rounded rect feel via border)."""
    rect(slide, l, t, w, h, CARD, border_color=color, border_pt=1.5)
    rect(slide, l, t, w, 0.055, color)
    txt = (icon + "  " if icon else "") + title
    label(slide, txt, l+0.15, t+0.1, w-0.25, 0.38, size=12, bold=True, color=color)
    if subtitle:
        label(slide, subtitle, l+0.15, t+0.52, w-0.25, h-0.6,
              size=10, color=MUTED)

def stat_card(slide, l, t, value, subtext, color):
    rect(slide, l, t, 3.1, 1.85, CARD)
    rect(slide, l, t, 0.06, 1.85, color)
    label(slide, value, l+0.2, t+0.2, 2.7, 0.75, size=34, bold=True, color=color)
    label(slide, subtext, l+0.2, t+0.95, 2.7, 0.72, size=11, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s1 = new_slide()
rect(s1, 0, 0, 13.33, 0.055, TEAL)
rect(s1, 0, 0, 0.08, 7.5, TEAL)
# decorative circle
sp = s1.shapes.add_shape(9, Inches(8.6), Inches(-1.0), Inches(6.5), Inches(6.5))
sp.fill.solid(); sp.fill.fore_color.rgb = CARD; sp.line.fill.background()

label(s1, "AI Engineering Internship  ·  April 2026",
      0.5, 0.28, 8.0, 0.4, size=11, color=MUTED)
label(s1, "Multi-Agent AI", 0.5, 1.0, 10.0, 1.3,
      size=60, bold=True, color=WHITE)
label(s1, "Academic Assistant", 0.5, 2.2, 10.5, 1.2,
      size=60, bold=True, color=TEAL)
rect(s1, 0.5, 3.55, 4.2, 0.045, TEAL)
label(s1, "Personalised tutoring · Real-time analytics · Student support",
      0.5, 3.7, 9.5, 0.5, size=15, color=LIGHT)
label(s1, "LangChain  ·  AutoGen  ·  GPT-4o-mini  ·  ChromaDB  ·  Streamlit",
      0.5, 4.25, 9.5, 0.45, size=12, color=MUTED)
for i, (t, col) in enumerate([("LangChain", TEAL), ("AutoGen", GREEN),
                                ("ChromaDB", PURPLE), ("Streamlit", AMBER)]):
    pill(s1, t, 0.5 + i * 2.05, 5.05, col)
label(s1, "Rachael Farhat", 0.5, 6.5, 5.0, 0.5, size=13, bold=True, color=WHITE)
label(s1, "Final Internship Project", 0.5, 6.95, 5.0, 0.38, size=11, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ═══════════════════════════════════════════════════════════════════════════════
s2 = new_slide()
slide_header(s2, "Problem Statement", "01 · Problem")

problems = [
    (TEAL,   "No Personalisation",
     "Most academic tools treat every student the same. They answer the question in front of them — but they have no idea what that student has struggled with before, which topics keep tripping them up, or how fast they actually learn."),
    (AMBER,  "No Feedback Loop",
     "Students often don't know where they're going wrong until an exam. There's no tool that automatically analyses their quiz history, spots their weak areas, and builds a study plan around those specific gaps."),
    (PURPLE, "No Holistic Support",
     "Academic help and emotional wellbeing are treated as completely separate things. But a student who is burnt out or overwhelmed needs more than another explanation — they need someone (or something) that recognises that."),
]
for i, (col, title, desc) in enumerate(problems):
    l = 0.5 + i * 4.25
    rect(s2, l, 1.3, 4.05, 5.3, CARD)
    rect(s2, l, 1.3, 4.05, 0.065, col)
    label(s2, title, l+0.2, 1.48, 3.65, 0.5, size=16, bold=True, color=col)
    label(s2, desc,  l+0.2, 2.1,  3.65, 4.2, size=12, color=LIGHT)

label(s2, "Single-agent chatbots answer one question at a time. They have no memory of who you are, no view of how you're doing, and no specialist depth across all three of these areas at once.",
      0.5, 6.75, 12.3, 0.6, size=11, color=MUTED, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Approach
# ═══════════════════════════════════════════════════════════════════════════════
s3 = new_slide()
slide_header(s3, "Approach: Four Specialised AI Agents", "02 · Approach")

label(s3, "Rather than one generalist chatbot, the system uses four agents — each built for a specific job. A central Orchestrator reads every incoming message, decides which specialist should handle it, and routes it accordingly. For complex queries that need more than one agent, it triggers a live multi-agent collaboration.",
      0.5, 1.18, 12.3, 0.72, size=12, color=MUTED)

agents = [
    (TEAL,   "🔀", "OrchestratorAgent",
     "The entry point for every query. It classifies intent into one of four categories — TUTOR, ANALYTICS, SUPPORT, or COLLABORATION — and routes the message to the right agent. For complex mixed queries, it launches an AutoGen group chat where multiple agents contribute to a single combined answer."),
    (GREEN,  "📚", "TutorAgent",
     "Handles all academic questions. It retrieves the most relevant chunks from the OpenStax course materials using RAG, then generates a grounded answer. It also produces adaptive multiple-choice quizzes on any topic, using structured output so the grading is always accurate."),
    (AMBER,  "📊", "AnalyticsAgent",
     "Tracks every quiz the student has taken in the current session. When asked about progress, it calculates their average score, identifies which topics they keep struggling with, and generates a personalised study plan targeting those exact gaps."),
    (PURPLE, "💬", "SupportAgent",
     "Handles the human side of studying. When a student says they're stressed, overwhelmed, or losing motivation, this agent responds with practical guidance on study strategies, time management, and wellbeing — rather than pushing more academic content at them."),
]
for i, (col, em, name, desc) in enumerate(agents):
    r = i // 2; c = i % 2
    l = 0.5 + c * 6.42; t = 2.05 + r * 2.55
    rect(s3, l, t, 6.15, 2.3, CARD)
    rect(s3, l, t, 0.06, 2.3, col)
    label(s3, em + "  " + name, l+0.2, t+0.12, 5.7, 0.45, size=14, bold=True, color=col)
    label(s3, desc, l+0.2, t+0.62, 5.7, 1.52, size=10.5, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — AI/ML Methods
# ═══════════════════════════════════════════════════════════════════════════════
s4 = new_slide()
slide_header(s4, "AI / ML Methods", "03 · Methods")

methods = [
    (TEAL,   "Large Language Model",
     "All agents are powered by OpenAI GPT-4o-mini via LangChain's ChatOpenAI. It handles everything from answering academic questions to classifying intent and generating quiz content."),
    (GREEN,  "Retrieval-Augmented Generation",
     "Instead of relying on the LLM's training data alone, the Tutor retrieves the five most relevant chunks from ChromaDB and injects them directly into the prompt — so every answer is grounded in actual course material."),
    (AMBER,  "Text Embeddings",
     "Each document chunk is converted to a 1536-dimensional vector using OpenAI's text-embedding-3-small model. These vectors are stored persistently in ChromaDB and searched by cosine similarity at query time."),
    (PURPLE, "Structured Output (Pydantic)",
     "Quiz questions are generated using ChatOpenAI.with_structured_output(), which returns a typed Python object with the correct answer and distractors as plain text. Python code then shuffles the options and assigns A/B/C/D — the LLM never picks a letter, so grading is always correct."),
    (TEAL,   "Multi-Agent Coordination",
     "For queries that need both tutoring knowledge and performance data, AutoGen's RoundRobinGroupChat runs TutorBot and AnalyticsBot in sequence. TextMentionTermination and a max turn limit keep the conversation from running indefinitely."),
    (GREEN,  "Zero-Shot Intent Classification",
     "The Orchestrator sends every query to the LLM with a system prompt listing the four valid categories. It reads the single-word reply and routes accordingly, falling back to TUTOR if the response is anything unexpected."),
]
for i, (col, title, desc) in enumerate(methods):
    r = i // 2; c = i % 2
    l = 0.5 + c * 6.42; t = 1.28 + r * 1.97
    rect(s4, l, t, 6.15, 1.78, CARD)
    rect(s4, l, t, 0.06, 1.78, col)
    label(s4, title, l+0.2, t+0.13, 5.7, 0.38, size=13, bold=True, color=col)
    label(s4, desc,  l+0.2, t+0.56, 5.7, 1.05, size=10.5, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Dataset & RAG Pipeline
# ═══════════════════════════════════════════════════════════════════════════════
s5 = new_slide()
slide_header(s5, "Dataset & RAG Pipeline", "04 · Data")

# Left panel — dataset
rect(s5, 0.5, 1.28, 5.9, 5.72, CARD)
rect(s5, 0.5, 1.28, 0.06, 5.72, TEAL)
label(s5, "OpenStax Textbooks  (CC BY 4.0)", 0.72, 1.4, 5.5, 0.45,
      size=14, bold=True, color=TEAL)
label(s5, "All course material comes from OpenStax — free, peer-reviewed university textbooks. Two books covering introductory CS and statistics form the core knowledge base.",
      0.72, 1.88, 5.45, 0.68, size=10.5, color=MUTED)
books = [
    ("Introduction to Computer Science",
     "Algorithms, data structures, programming fundamentals"),
    ("Introductory Statistics, 2e",
     "Probability, hypothesis testing, regression analysis"),
]
for i, (title, sub) in enumerate(books):
    t = 2.72 + i * 1.38
    rect(s5, 0.72, t, 5.5, 1.15, RGBColor(0x0F, 0x1C, 0x35))
    rect(s5, 0.72, t, 0.045, 1.15, GREEN)
    label(s5, title, 0.88, t+0.1, 5.15, 0.38, size=12, bold=True, color=WHITE)
    label(s5, sub,   0.88, t+0.5, 5.15, 0.45, size=10.5, color=MUTED)
label(s5, "Students can also upload their own PDFs, DOCX files, or plain text at any time — these are embedded on the fly and become immediately searchable within that session.",
      0.72, 5.6, 5.45, 0.75, size=10.5, color=MUTED, italic=True)

# Right panel — pipeline
label(s5, "How RAG Works", 6.75, 1.28, 6.1, 0.45, size=14, bold=True, color=TEAL)
steps = [
    (1, "Load",     "PyPDFLoader reads every PDF in /data/course_catalogue/",            TEAL),
    (2, "Chunk",    "Split into 800-character pieces with 100-char overlap",              GREEN),
    (3, "Embed",    "text-embedding-3-small converts each chunk to a 1536-dim vector",   AMBER),
    (4, "Store",    "Vectors saved to ChromaDB on disk with a source metadata tag",      PURPLE),
    (5, "Retrieve", "At query time, the five most similar chunks are fetched",            TEAL),
    (6, "Answer",   "Those chunks are injected into the LLM prompt as grounding context",GREEN),
]
for i, (n, title, detail, col) in enumerate(steps):
    t = 1.88 + i * 0.9
    rect(s5, 6.75, t, 0.52, 0.52, col)
    label(s5, str(n), 6.9, t+0.07, 0.28, 0.35, size=16, bold=True, color=NAVY)
    label(s5, title,  7.38, t+0.04, 5.85, 0.3,  size=12, bold=True, color=WHITE)
    label(s5, detail, 7.38, t+0.32, 5.85, 0.35, size=10, color=MUTED)
    if i < len(steps) - 1:
        rect(s5, 6.97, t+0.54, 0.06, 0.36, MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — System Architecture (proper block diagram)
# ═══════════════════════════════════════════════════════════════════════════════
s6 = new_slide()
slide_header(s6, "System Architecture", "05 · Architecture")

# ── Layout (all measurements in Inches) ──────────────────────────────────────
# Row 1 (y=1.3):  [Student]  →  [Streamlit UI]  →  [OrchestratorAgent]
# Row 2 (y=3.1):  [TutorAgent]  [AnalyticsAgent]  [SupportAgent]
#                 ↑ all three connected from Orchestrator above
# Row 3 (y=5.1):  [ChromaDB]  (under Tutor)     [AutoGen] (right, from Orchestrator)

NODE_H = 1.0

# Row 1
node(s6, 0.4,  1.3,  2.3, NODE_H, "Student",              "Sends a question or\ntakes a quiz",                     MUTED,  "👤")
node(s6, 3.4,  1.3,  2.8, NODE_H, "Streamlit UI",          "Chat · Quiz · Analytics · Upload",                      TEAL,   "🖥")
node(s6, 7.0,  1.3,  5.8, NODE_H, "OrchestratorAgent",     "Classifies intent → routes to specialist\nor starts AutoGen group chat", AMBER, "🔀")

# H-arrows Row 1
arrow_h(s6, 2.72, 1.72, 0.68, TEAL)
arrow_h(s6, 6.22, 1.72, 0.78, AMBER)

# Row 2 — three specialist agents, centred under the orchestrator span
node(s6, 0.4,  3.1,  3.75, NODE_H, "TutorAgent",      "RAG Q&A grounded in course\nmaterials + quiz generation",  GREEN,  "📚")
node(s6, 4.65, 3.1,  3.75, NODE_H, "AnalyticsAgent",  "Quiz history · weak area\ndetection · study plans",        PURPLE, "📊")
node(s6, 8.9,  3.1,  3.9,  NODE_H, "SupportAgent",    "Wellbeing · motivation\n· study strategy",                 TEAL,   "💬")

# V-arrows from Orchestrator bottom to each agent top
# Orchestrator occupies x=7.0 to 12.8, centre ≈ 9.9
# We draw drop lines from the mid-bottom of orchestrator to a horizontal bus,
# then branch to each agent.
orch_cx = 9.9   # horizontal centre of orchestrator box
bus_y   = 2.5   # y of the horizontal bus line
arrow_v(s6, orch_cx - 0.03, 2.3, bus_y - 2.3 - 0.1, AMBER)

# Horizontal bus across all three agents
rect(s6, 2.25, bus_y, 8.6, 0.04, MUTED)

# Vertical drops from bus to each agent
for cx in [2.28, 6.53, 10.85]:
    arrow_v(s6, cx - 0.03, bus_y, 3.1 - bus_y - 0.06, MUTED)

# Row 3 — ChromaDB (under Tutor) and AutoGen (beside Orchestrator area)
node(s6, 0.4,  4.8,  3.75, 1.05, "ChromaDB",
     "Persistent vector store\nOpenStax chunks + user uploads", AMBER, "🗄")
node(s6, 4.65, 4.8,  3.75, 1.05, "AutoGen Group Chat",
     "TutorBot + AnalyticsBot collaborate\non COLLABORATION-intent queries", GREEN, "⚡")
node(s6, 8.9,  4.8,  3.9,  1.05, "DuckDuckGo Search",
     "Optional toggle — live web\ncontext beyond course materials", PURPLE, "🔍")

# V-arrows from agents to row 3
arrow_v(s6, 2.25,  4.1, 0.68, GREEN)
arrow_v(s6, 6.5,   4.1, 0.68, GREEN)

# Label annotation
label(s6, "* OrchestratorAgent connects directly to AutoGen for COLLABORATION queries; "
          "TutorAgent pulls from ChromaDB for every RAG-grounded response.",
      0.4, 6.15, 12.5, 0.55, size=9.5, color=MUTED, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Features (screenshots) — Part 1: Chat & Analytics
# ═══════════════════════════════════════════════════════════════════════════════
s7 = new_slide()
slide_header(s7, "Features — Chat Tutor & Analytics", "06 · Features")

# Left: Chat
rect(s7, 0.5, 1.2, 6.1, 4.6, CARD, border_color=TEAL, border_pt=1)
tf = s7.shapes[-1].text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run(); run.text = "[ Screenshot: Chat + RAG Tutor ]"
run.font.size = Pt(11); run.font.color.rgb = MUTED; run.font.italic = True

label(s7, "Chat Tutor", 0.5, 5.88, 3.0, 0.38, size=13, bold=True, color=TEAL)
label(s7, "Ask any academic question and the TutorAgent retrieves the most relevant passages from the OpenStax textbooks using RAG, then constructs a grounded answer. You can also toggle on live web search for context that goes beyond the course materials.",
      0.5, 6.28, 6.1, 0.95, size=10.5, color=MUTED)

# Right: Analytics
rect(s7, 6.95, 1.2, 6.1, 4.6, CARD, border_color=GREEN, border_pt=1)
tf2 = s7.shapes[-1].text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run(); run2.text = "[ Screenshot: Analytics Dashboard ]"
run2.font.size = Pt(11); run2.font.color.rgb = MUTED; run2.font.italic = True

label(s7, "Analytics Dashboard", 6.95, 5.88, 5.0, 0.38, size=13, bold=True, color=GREEN)
label(s7, "A Plotly bar chart shows quiz scores per topic, colour-coded by performance tier — green for 80%+, amber for 60–79%, red below 60%. Below the chart, the AnalyticsAgent generates a study plan targeting the student's weakest areas.",
      6.95, 6.28, 6.1, 0.95, size=10.5, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Features (screenshots) — Part 2: Quiz & Upload
# ═══════════════════════════════════════════════════════════════════════════════
s8 = new_slide()
slide_header(s8, "Features — Quiz & Document Upload", "07 · Features")

# Left: Quiz
rect(s8, 0.5, 1.2, 6.1, 4.6, CARD, border_color=AMBER, border_pt=1)
tf3 = s8.shapes[-1].text_frame; tf3.word_wrap = True
p3 = tf3.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
run3 = p3.add_run(); run3.text = "[ Screenshot: Quiz Generator ]"
run3.font.size = Pt(11); run3.font.color.rgb = MUTED; run3.font.italic = True

label(s8, "Quiz Generator", 0.5, 5.88, 3.0, 0.38, size=13, bold=True, color=AMBER)
label(s8, "Enter any topic and the system generates multiple-choice questions grounded in course material. Pydantic structured output guarantees that the correct answer is always marked correctly — the LLM returns text only, and Python assigns the option letters after shuffling.",
      0.5, 6.28, 6.1, 0.95, size=10.5, color=MUTED)

# Right: Upload / Summariser
rect(s8, 6.95, 1.2, 6.1, 4.6, CARD, border_color=PURPLE, border_pt=1)
tf4 = s8.shapes[-1].text_frame; tf4.word_wrap = True
p4 = tf4.paragraphs[0]; p4.alignment = PP_ALIGN.CENTER
run4 = p4.add_run(); run4.text = "[ Screenshot: Document Upload & Summariser ]"
run4.font.size = Pt(11); run4.font.color.rgb = MUTED; run4.font.italic = True

label(s8, "Document Upload & Summariser", 6.95, 5.88, 5.5, 0.38, size=13, bold=True, color=PURPLE)
label(s8, "Students can upload their own PDF, DOCX, or TXT files. The document is chunked and embedded immediately, making it searchable in the same RAG pipeline alongside the course catalogue. A one-click summariser produces a structured overview of any uploaded material.",
      6.95, 6.28, 6.1, 0.95, size=10.5, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Results
# ═══════════════════════════════════════════════════════════════════════════════
s9 = new_slide()
slide_header(s9, "Results", "08 · Results")

stats = [
    ("4",      "Specialist AI agents\nfully operational",         TEAL),
    ("~3 K",   "OpenStax pages indexed\nin ChromaDB",             GREEN),
    ("100%",   "Quiz grading accuracy\nvia Pydantic output",      AMBER),
    ("5",      "Core features shipped\n(Chat · Quiz · Analytics\n·Upload · Support)", PURPLE),
]
for i, (val, sub, col) in enumerate(stats):
    stat_card(s9, 0.5 + i * 3.21, 1.28, val, sub, col)

label(s9, "What was achieved:", 0.5, 3.42, 5.0, 0.38, size=13, bold=True, color=TEAL)
outcomes = [
    ("All four agents are operational",
     "Intent routing has been tested and verified across tutoring, analytics, support, and collaboration query types."),
    ("RAG delivers accurate, grounded answers",
     "ChromaDB retrieves relevant chunks from the OpenStax textbooks, keeping responses tied to actual course content rather than LLM guesswork."),
    ("Quiz grading is now fully reliable",
     "Switching to Pydantic structured output eliminated the persistent grading bug — the LLM no longer assigns option letters, so answers can never be mislabelled."),
    ("AutoGen collaboration works end-to-end",
     "For complex queries that need both tutoring depth and personalised analytics, TutorBot and AnalyticsBot collaborate and return a single combined response."),
]
for i, (title, detail) in enumerate(outcomes):
    t = 3.95 + i * 0.82
    rect(s9, 0.5, t, 0.18, 0.18, TEAL)
    label(s9, title,  0.82, t-0.04, 12.0, 0.3, size=12, bold=True, color=WHITE)
    label(s9, detail, 0.82, t+0.28, 12.0, 0.4, size=11, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Visualizations
# ═══════════════════════════════════════════════════════════════════════════════
s10 = new_slide()
slide_header(s10, "Visualizations", "09 · Visualizations")

label(s10, "The analytics dashboard is built with Plotly and gives students a clear view of where they stand across every topic they've been quizzed on.",
      0.5, 1.18, 12.3, 0.48, size=12, color=MUTED)

# Main screenshot placeholder
rect(s10, 0.5, 1.75, 7.85, 4.2, CARD, border_color=TEAL, border_pt=1)
sp = s10.shapes[-1]
tf = sp.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run(); run.text = "[ Screenshot: Plotly analytics bar chart ]"
run.font.size = Pt(11); run.font.color.rgb = MUTED; run.font.italic = True

# Side callout cards
side = [
    (TEAL,   "Colour-coded bars",        "Green ≥ 80% · Amber ≥ 60% · Red < 60%"),
    (GREEN,  "Personalised study plan",  "AnalyticsAgent text output appears below the chart with targeted recommendations"),
    (AMBER,  "Session stats in sidebar", "Total questions answered, current average score, and streak counter"),
    (PURPLE, "Instant quiz feedback",    "Each answer shows whether it was correct and why, immediately after submission"),
]
for i, (col, title, desc) in enumerate(side):
    t = 1.75 + i * 1.07
    rect(s10, 8.6, t, 4.45, 0.9, CARD)
    rect(s10, 8.6, t, 0.055, 0.9, col)
    label(s10, title, 8.78, t+0.08, 4.1, 0.32, size=12, bold=True, color=col)
    label(s10, desc,  8.78, t+0.44, 4.1, 0.38, size=10, color=MUTED)

rect(s10, 0.5, 6.05, 7.85, 0.9, CARD, border_color=GREEN, border_pt=1)
sp2 = s10.shapes[-1]
tf2 = sp2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run(); run2.text = "[ Screenshot: AnalyticsAgent study plan output ]"
run2.font.size = Pt(10); run2.font.color.rgb = MUTED; run2.font.italic = True


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Challenges & Solutions
# ═══════════════════════════════════════════════════════════════════════════════
s11 = new_slide()
slide_header(s11, "Challenges & How We Solved Them", "10 · Challenges")

challenges = [
    (TEAL,   "Quiz grading always wrong",
     "No matter how the prompt was worded, the LLM kept assigning the wrong letter to the correct answer.",
     "Switched to Pydantic structured output — the LLM now returns the correct answer and distractors as plain text strings. Python shuffles the four options and assigns A/B/C/D in code. The LLM never touches the letters, so grading cannot go wrong."),
    (AMBER,  "LaTeX not rendering in the UI",
     "Mathematical expressions were appearing as raw \\frac{}{} strings rather than formatted equations.",
     "Added a post-processing step that converts all LLM output to KaTeX-compatible $...$ notation before it's passed to Streamlit's markdown renderer."),
    (GREEN,  "Stale agent instances after code changes",
     "After updating agent logic, Streamlit's session state kept serving the old cached OrchestratorAgent object.",
     "Added a _version class attribute. The get_orchestrator() function compares the cached object's version against the current class version and rebuilds if they don't match."),
    (PURPLE, "App was slow to start",
     "ChromaDB and the OpenAI embedding model were being reinitialised on every single page interaction.",
     "Wrapped both vector stores with @st.cache_resource so they're only created once per server process. Catalogue indexing was also made lazy — it runs on the first query, not at startup."),
    (TEAL,   "Wrong files appearing in the catalogue",
     "iterdir() was picking up hidden files and OS metadata files alongside the actual PDFs.",
     "Added an explicit suffix check so only .pdf, .docx, and .txt files are treated as course materials."),
]
for i, (col, title, prob, fix) in enumerate(challenges):
    r = i // 2; c = i % 2
    if i == 4:
        l, t, w = 0.5, 1.28 + r * 2.52, 12.45
    else:
        l, t, w = 0.5 + c * 6.47, 1.28 + r * 2.52, 6.2
    rect(s11, l, t, w, 2.22, CARD)
    rect(s11, l, t, 0.06, 2.22, col)
    label(s11, title,         l+0.2, t+0.1,  w-0.35, 0.38, size=13, bold=True, color=col)
    label(s11, "Problem: " + prob, l+0.2, t+0.54, w-0.35, 0.6,  size=10, color=MUTED)
    label(s11, "Fix: " + fix,      l+0.2, t+1.14, w-0.35, 0.9,  size=11, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Key Learnings
# ═══════════════════════════════════════════════════════════════════════════════
s12 = new_slide()
slide_header(s12, "Key Learnings", "11 · Learnings")

learnings = [
    (TEAL,   "LLMs need clear instructions to route correctly",
     "Without telling the model exactly which categories exist and what falls into each one, it will sometimes handle queries it should have passed to a specialist. A simple system prompt with four named categories fixed this completely."),
    (GREEN,  "Structured output is the right tool for precise tasks",
     "Asking the LLM to generate quiz questions AND assign the correct letter in free text always broke. Separating the two — LLM writes the content, Python handles the structure — made the quiz reliable immediately."),
    (AMBER,  "How you split documents affects answer quality",
     "Chunks that were too large pulled in unrelated content. Chunks that were too small cut off explanations halfway through. Finding the right size (800 characters, 100 overlap) made a noticeable difference in answer accuracy."),
    (PURPLE, "Caching is essential for AI apps to feel fast",
     "Without caching, the app was re-loading the entire document library on every click — taking around 20 seconds each time. One line of Streamlit caching brought that down to under 2 seconds."),
    (TEAL,   "Multi-agent systems need a defined stopping point",
     "The first time AutoGen ran without a turn limit, it just kept going. You have to explicitly tell the system when a conversation is finished, or it won't stop on its own."),
    (GREEN,  "Real content beats made-up examples",
     "Using actual OpenStax textbooks meant the assistant was answering questions from real course material from day one. Students got accurate, relevant answers rather than generic explanations the model invented."),
]
for i, (col, title, desc) in enumerate(learnings):
    r = i // 2; c = i % 2
    l = 0.5 + c * 6.47; t = 1.28 + r * 2.0
    rect(s12, l, t, 6.2, 1.8, CARD)
    rect(s12, l, t, 0.06, 1.8, col)
    label(s12, title, l+0.2, t+0.12, 5.8, 0.42, size=13, bold=True, color=col)
    label(s12, desc,  l+0.2, t+0.6,  5.8, 1.05, size=10.5, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Closing
# ═══════════════════════════════════════════════════════════════════════════════
s13 = new_slide()
rect(s13, 0, 0, 13.33, 0.055, TEAL)
rect(s13, 0, 0, 0.08, 7.5, TEAL)
sp = s13.shapes.add_shape(9, Inches(7.8), Inches(-0.8), Inches(6.5), Inches(6.5))
sp.fill.solid(); sp.fill.fore_color.rgb = CARD; sp.line.fill.background()

label(s13, "Thank You", 0.5, 1.4, 9.5, 1.6, size=66, bold=True, color=WHITE)
label(s13, "Multi-Agent AI Academic Assistant", 0.5, 3.0, 9.5, 0.65,
      size=20, bold=True, color=TEAL)
rect(s13, 0.5, 3.78, 3.5, 0.045, TEAL)
label(s13, "Four specialised agents. One seamless experience for students.",
      0.5, 3.95, 9.5, 0.48, size=13, color=MUTED, italic=True)

for i, (t, col) in enumerate([("4 AI Agents", TEAL), ("RAG · ~3K pages", GREEN),
                               ("100% quiz accuracy", AMBER), ("LMS-style UI", PURPLE)]):
    pill(s13, t, 0.5 + i * 3.1, 4.65, col, size=12)

label(s13, "Questions?", 0.5, 5.4, 9.0, 0.7, size=30, bold=True, color=TEAL)
label(s13, "Rachael Farhat  ·  AI Engineering Internship  ·  April 2026",
      0.5, 6.3, 9.0, 0.45, size=12, color=MUTED)


prs.save("presentation_final.pptx")
print("Saved: presentation_final.pptx  (13 slides, 16:9)")
print("Replace screenshot placeholders in slides 7, 8, and 10 with actual app screenshots.")
